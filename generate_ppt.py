"""Generate 30-slide academic PPT."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "results")

DARK = RGBColor(0x1a,0x1a,0x2e); ACCENT = RGBColor(0x1d,0x4e,0xd8)
WHITE = RGBColor(0xFF,0xFF,0xFF); LIGHT = RGBColor(0xF0,0xF4,0xFF)
GRAY = RGBColor(0x64,0x74,0x8B); GREEN = RGBColor(0x16,0xa3,0x4a)
RED = RGBColor(0xdc,0x26,0x26)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s, c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def bar(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def txt(s, l, t, w, h, text, sz=18, b=False, c=DARK, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
    p.font.name = "Calibri"; p.alignment = a
    return tb.text_frame

def bullet_slide(title, items, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    bar(s, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)
    txt(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6), title, 30, True)
    bar(s, Inches(0.8), Inches(0.95), Inches(1.8), Inches(0.04), ACCENT)
    y = 1.2
    if sub:
        txt(s, Inches(0.8), Inches(y), Inches(11), Inches(0.4), sub, 15, c=GRAY); y = 1.6
    tb = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(7.5-y-0.3))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.name = "Calibri"; p.space_after = Pt(5)
        if item.startswith("  "):
            p.level = 1; p.font.size = Pt(16); p.font.color.rgb = GRAY
        elif item == "":
            p.font.size = Pt(8)
        else:
            p.font.size = Pt(18); p.font.color.rgb = DARK
    return s

def img_slide(title, img, caption=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    bar(s, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)
    txt(s, Inches(0.8), Inches(0.25), Inches(11), Inches(0.5), title, 26, True)
    p = os.path.join(RESULTS, img)
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(1.5), Inches(1.0), Inches(10), Inches(5.6))
    if caption:
        txt(s, Inches(1), Inches(6.7), Inches(11), Inches(0.5), caption, 13, c=GRAY, a=PP_ALIGN.CENTER)

def tbl_slide(title, headers, rows, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    bar(s, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)
    txt(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6), title, 30, True)
    bar(s, Inches(0.8), Inches(0.95), Inches(1.8), Inches(0.04), ACCENT)
    t = s.shapes.add_table(len(rows)+1, len(headers), Inches(2), Inches(1.4), Inches(9), Inches(0.5+len(rows)*0.55)).table
    for i, h in enumerate(headers):
        c = t.cell(0, i); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold=True; p.font.size=Pt(16); p.font.color.rgb=WHITE; p.alignment=PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for r, row in enumerate(rows):
        for ci, v in enumerate(row):
            c = t.cell(r+1, ci); c.text = v
            for p in c.text_frame.paragraphs:
                p.font.size=Pt(15); p.alignment=PP_ALIGN.CENTER
                if ci > 0 and any(x in v for x in ["−","97","98"]): p.font.bold=True
            if r % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LIGHT
    if note:
        txt(s, Inches(2), Inches(1.6+len(rows)*0.55+0.8), Inches(9), Inches(0.6), note, 14, c=GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
bar(s, Inches(0), Inches(2.8), prs.slide_width, Inches(0.05), ACCENT)
txt(s, Inches(1), Inches(1.2), Inches(11), Inches(1.2), "Automated Industrial Environmental\nCompliance Monitoring System", 38, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.1), Inches(11), Inches(0.6), "Using Multi-Temporal Sentinel-2 Imagery and Machine Learning", 22, c=RGBColor(0x93,0xAE,0xDB), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5), "Oishik Kar  \u2022  Mihir Dagar  \u2022  Shouvik Das  \u2022  Youvansh Chauhan", 18, c=RGBColor(0xA0,0xA0,0xC0), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.1), Inches(11), Inches(0.5), "Under the guidance of Dr. Sunil C K, Assistant Professor", 16, c=RGBColor(0x80,0x80,0xA0), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.7), Inches(11), Inches(0.5), "Department of CSE, IIIT Dharwad  |  ML Course Project \u2014 Statement 3", 15, c=RGBColor(0x70,0x70,0x90), a=PP_ALIGN.CENTER)

# === SLIDES 2-17: Content ===
bullet_slide("Problem Statement", [
    "Manual environmental inspections of industrial zones are infrequent and hard to scale",
    "Karnataka SPCB requires factories to maintain minimum green cover \u2014 compliance rarely verified",
    "Need: Automated system using satellite imagery + ML to detect violations continuously",
    "", "Statement 3 Requirements:",
    "  \u2022 Utilize KGIS boundary data to identify factory premises",
    "  \u2022 Analyze satellite imagery over time \u2014 detect green cover loss & unauthorized construction",
    "  \u2022 Apply ML models for vegetation detection and land-use change classification",
    "  \u2022 Generate compliance reports with geo-coordinates, temporal analysis, visual proof",
    "  \u2022 Enable scalable, continuous monitoring for pollution control boards",
])

bullet_slide("Objectives", [
    "1. Build an end-to-end automated pipeline from data acquisition to compliance reporting",
    "2. Integrate KGIS taluk boundary data for precise industrial zone delineation",
    "3. Download and preprocess multi-temporal Sentinel-2 imagery (2020 vs 2024)",
    "4. Compute spectral indices (NDVI, NDBI, NBI) for vegetation and built-up analysis",
    "5. Train Random Forest classifier for 6-class land cover classification",
    "6. Perform bitemporal change detection to quantify green cover loss",
    "7. Extract geo-tagged violation zones with coordinates and area estimates",
    "8. Generate HTML compliance reports and an interactive web dashboard",
])

bullet_slide("Study Area \u2014 Peenya Industrial Area", [
    "Location: North-west Bengaluru, Karnataka",
    "KGIS Taluk: Bangalore-North (Code 2001)",
    "Coordinates: 13.01\u201313.07\u00b0N, 77.48\u201377.56\u00b0E",
    "Area: ~40 km\u00b2",
    "", "Characteristics:",
    "  \u2022 One of the largest industrial estates in South Asia",
    "  \u2022 8,000+ registered manufacturing, electronics, and heavy industry units",
    "  \u2022 Significant land-use transition due to metropolitan expansion",
    "  \u2022 Dense mix of factories with mandated green belt requirements",
])

bullet_slide("Study Area \u2014 Whitefield Industrial Area", [
    "Location: East Bengaluru, Karnataka",
    "KGIS Taluk: Bangalore-East (Code 2004)",
    "Coordinates: 12.95\u201313.01\u00b0N, 77.72\u201377.80\u00b0E",
    "Area: ~25 km\u00b2",
    "", "Characteristics:",
    "  \u2022 Rapidly urbanizing IT corridor (ITPL, major tech parks)",
    "  \u2022 Mixed industrial and residential development",
    "  \u2022 Heavy construction activity driving green cover loss",
    "  \u2022 Different industrial profile from Peenya \u2014 enables comparison",
])

bullet_slide("Data Source 1 \u2014 KGIS (Karnataka GIS)", [
    "Source: kgis.ksrsac.in (KSRSAC portal)",
    "What it provides:",
    "  \u2022 Administrative boundary shapefiles (taluk, district, ward level)",
    "  \u2022 Downloaded: Taluk.shp with all 226 Karnataka taluks",
    "  \u2022 Filtered by: Code 2001 (Bangalore-North) and 2004 (Bangalore-East)",
    "  \u2022 Location Details REST API for ward/zone metadata per violation",
    "", "Purpose in project:",
    "  \u2022 Delineates industrial zone boundaries \u2014 tells system WHERE to analyze",
    "  \u2022 Masks satellite imagery to only include pixels inside industrial premises",
    "", "Does NOT provide: satellite imagery, vegetation data, or temporal information",
])

bullet_slide("Data Source 2 \u2014 Sentinel-2 L2A Satellite Imagery", [
    "Source: ESA Copernicus via Microsoft Planetary Computer (STAC API)",
    "Free, open-access multispectral satellite imagery",
    "", "Specifications:",
    "  \u2022 Resolution: 10m (B02\u2013B08), 20m resampled to 10m (B11\u2013B12)",
    "  \u2022 Bands: Blue, Green, Red, NIR, SWIR1, SWIR2 + SCL cloud mask",
    "  \u2022 Time periods: March 2020 (baseline) vs March 2024 (recent)",
    "  \u2022 Cloud cover: <0.2% (T1), <0.01% (T2) \u2014 dry season selection",
    "", "Scale:",
    "  \u2022 4 satellite scenes across 28 GeoTIFF files",
    "  \u2022 ~11.6 million multispectral pixels total",
    "  \u2022 582,066 training + 249,458 test samples for ML",
])

bullet_slide("System Architecture \u2014 Stages 1\u20133", [
    "Stage 1: Data Acquisition",
    "  \u2022 KGIS shapefile download from KSRSAC portal",
    "  \u2022 Sentinel-2 STAC API query with cloud cover filter (<15%)",
    "  \u2022 Windowed GeoTIFF read (downloads only AOI, not full 100km tile)",
    "", "Stage 2: Preprocessing",
    "  \u2022 SCL cloud masking (removes cloud, shadow, cirrus pixels)",
    "  \u2022 20m \u2192 10m bilinear resampling (B11, B12)",
    "  \u2022 KGIS boundary rasterization and spatial masking",
    "  \u2022 7-band GeoTIFF stacking [B02, B03, B04, B08, B11, B12, Mask]",
    "", "Stage 3: Spectral Index Computation",
    "  \u2022 NDVI: vegetation health quantification",
    "  \u2022 NDBI: built-up area detection",
    "  \u2022 NBI: new construction identification",
])

bullet_slide("System Architecture \u2014 Stages 4\u20136", [
    "Stage 4: ML Classification",
    "  \u2022 Random Forest (100 trees) trained on 6 raw spectral bands",
    "  \u2022 K-Means (k=5) unsupervised clustering as baseline",
    "  \u2022 70/30 stratified train-test split with validation metrics",
    "", "Stage 5: Change Detection",
    "  \u2022 Bitemporal NDVI/NDBI differencing (2024 \u2212 2020)",
    "  \u2022 Morphological closing to merge scattered detections",
    "  \u2022 Connected component labeling for violation zone extraction",
    "", "Stage 6: Report Generation",
    "  \u2022 HTML compliance reports with risk assessment",
    "  \u2022 Annotated satellite imagery with violation markers",
    "  \u2022 Geo-coordinate export (lat/lon for each violation)",
    "  \u2022 Interactive web dashboard with Leaflet.js maps",
])

bullet_slide("Preprocessing \u2014 Cloud Masking & Resampling", [
    "Cloud Masking (SCL Band)",
    "  \u2022 Scene Classification Layer identifies pixel types",
    "  \u2022 Masked: Cloud shadow (3), Cloud medium (8), Cloud high (9), Cirrus (10)",
    "  \u2022 Valid: Vegetation (4), Not-vegetated (5), Water (6), Unclassified (7)",
    "  \u2022 Result: 99.8% valid pixel coverage across all 4 acquisitions",
    "", "Band Resampling",
    "  \u2022 B11 (SWIR1) and B12 (SWIR2) at 20m native resolution",
    "  \u2022 Resampled to 10m using bilinear interpolation (scipy.ndimage.zoom)",
    "", "KGIS Boundary Masking",
    "  \u2022 Taluk polygon reprojected from WGS84 to UTM Zone 43N",
    "  \u2022 Rasterized onto satellite grid \u2014 pixels outside boundary = 0",
    "  \u2022 Combined mask: pixel must be cloud-free AND inside KGIS boundary",
])

bullet_slide("Spectral Indices", [
    "NDVI \u2014 Normalized Difference Vegetation Index",
    "  \u2022 Formula: (NIR \u2212 Red) / (NIR + Red)",
    "  \u2022 Range: [\u22121, +1] \u2014 higher values = more vegetation",
    "  \u2022 Used for: Green cover quantification and vegetation loss detection",
    "", "NDBI \u2014 Normalized Difference Built-up Index",
    "  \u2022 Formula: (SWIR1 \u2212 NIR) / (SWIR1 + NIR)",
    "  \u2022 Range: [\u22121, +1] \u2014 higher values = more concrete/built-up",
    "  \u2022 Used for: Unauthorized construction detection",
    "", "NBI \u2014 New Built-up Index",
    "  \u2022 Formula: (Red \u00d7 SWIR1) / NIR",
    "  \u2022 Higher values = bare soil or fresh construction",
])

bullet_slide("Change Detection Methodology", [
    "Bitemporal Differencing",
    "  \u2022 \u0394NDVI = NDVI_2024 \u2212 NDVI_2020 \u2192 negative = vegetation lost",
    "  \u2022 \u0394NDBI = NDBI_2024 \u2212 NDBI_2020 \u2192 positive = new construction",
    "", "Classification Thresholds",
    "  \u2022 Class 1 \u2014 Vegetation Loss: \u0394NDVI < \u22120.15",
    "  \u2022 Class 2 \u2014 New Construction: \u0394NDBI > +0.10",
    "  \u2022 Class 3 \u2014 Both conditions met simultaneously",
    "  \u2022 Class 0 \u2014 No significant change",
    "", "Risk Assessment",
    "  \u2022 HIGH: green cover change < \u221230%",
    "  \u2022 MEDIUM: green cover change between \u221215% and \u221230%",
    "  \u2022 LOW: green cover change > \u221215%",
])

bullet_slide("ML Classification \u2014 Random Forest", [
    "Algorithm: Random Forest (scikit-learn)",
    "  \u2022 100 decision trees, max_depth=15, balanced class weights",
    "  \u2022 min_samples_split=10, min_samples_leaf=5, n_jobs=-1",
    "", "Features (input): 6 raw spectral bands",
    "  \u2022 B02 (Blue), B03 (Green), B04 (Red), B08 (NIR), B11 (SWIR1), B12 (SWIR2)",
    "", "Labels (output): 6 land cover classes",
    "  \u2022 Dense Vegetation, Sparse Vegetation, Built-up, Bare Soil, Water, Other",
    "", "Training setup:",
    "  \u2022 Each 10m\u00d710m pixel = one sample with 6 reflectance values",
    "  \u2022 70% train / 30% test, stratified by class",
    "  \u2022 K-Means (k=5) as unsupervised baseline for comparison",
])

bullet_slide("Non-Circular Validation \u2014 Key Design Decision", [
    "The Problem:",
    "  \u2022 Training labels are derived from NDVI/NDBI thresholds",
    "  \u2022 If RF trains on NDVI/NDBI as features \u2192 trivially memorizes thresholds",
    "  \u2022 Would get meaningless 100% accuracy",
    "", "Our Solution:",
    "  \u2022 Labels generated from: NDVI and NDBI (spectral indices)",
    "  \u2022 RF trained on: Raw bands ONLY (B02, B03, B04, B08, B11, B12)",
    "  \u2022 Features and label sources are completely DISJOINT",
    "", "Why This Works:",
    "  \u2022 RF must learn spectral reflectance patterns from raw data",
    "  \u2022 Cannot shortcut by reading NDVI/NDBI values directly",
    "  \u2022 The 97\u201398% accuracy reflects genuine pattern recognition",
])

bullet_slide("Violation Extraction & Geo-Tagging", [
    "Morphological Processing",
    "  \u2022 5\u00d75 structuring element, 3 iterations of closing",
    "  \u2022 Merges scattered single-pixel detections within 50m into regions",
    "", "Connected Component Labeling",
    "  \u2022 SciPy labels contiguous violation clusters",
    "  \u2022 Filters regions below 500 m\u00b2 (noise removal)",
    "", "Geo-Coordinate Conversion",
    "  \u2022 Centroid pixel \u2192 UTM Zone 43N \u2192 WGS84 (lat/lon) via PyProj",
    "  \u2022 Each violation: latitude, longitude, type, area (hectares)",
    "", "KGIS API Enrichment",
    "  \u2022 Queries KGIS Location Details REST API per violation",
    "  \u2022 Adds ward name, zone name, town, district metadata",
])

# === SLIDES 18-19: Tables ===
tbl_slide("Results \u2014 Change Detection", ["Metric", "Peenya", "Whitefield"], [
    ["Green Cover 2020", "52.57%", "71.25%"],
    ["Green Cover 2024", "28.51%", "43.45%"],
    ["Net Change", "\u221224.06%", "\u221227.80%"],
    ["Vegetation Loss", "27.23%", "38.92%"],
    ["New Construction", "4.56%", "17.72%"],
    ["Violation Zones", "54", "72"],
    ["Risk Level", "MEDIUM", "MEDIUM"],
])

tbl_slide("Results \u2014 ML Model Performance", ["Metric", "Peenya", "Whitefield"], [
    ["Training Samples", "405,333", "176,733"],
    ["Test Samples", "173,715", "75,743"],
    ["Accuracy", "97.74%", "98.70%"],
    ["Precision (weighted)", "98.06%", "98.79%"],
    ["Recall (weighted)", "97.74%", "98.70%"],
    ["F1-Score (weighted)", "97.83%", "98.73%"],
], "Non-circular validation: RF trained on raw bands only, labels from NDVI/NDBI thresholds")

# === SLIDES 20-26: Images ===
for t, f, c in [
    ("Results \u2014 Confusion Matrix (Peenya)", "confusion_matrix_peenya.png", "6\u00d76 matrix on 173,715 test pixels \u2014 strong diagonal dominance across all classes"),
    ("Results \u2014 Feature Importance", "feature_importance_peenya.png", "B04 (Red, 28.4%) and B08 (NIR, 23.7%) dominate \u2014 consistent with vegetation spectral theory"),
    ("Results \u2014 Learning Curve", "learning_curve_peenya.png", "Accuracy stabilizes at ~50 trees, confirming adequate model capacity without overfitting"),
    ("Results \u2014 Before/After Analysis (Peenya)", "before_after_peenya.png", "Top: RGB 2020 vs 2024 | Bottom: NDVI heatmaps showing vegetation density decline"),
    ("Results \u2014 Before/After Analysis (Whitefield)", "before_after_whitefield.png", "Severe vegetation loss in IT corridor \u2014 71.25% \u2192 43.45% green cover"),
    ("Results \u2014 Annotated Violation Map", "annotated_violations_whitefield.png", "Red: vegetation loss | Orange: new construction | Pink: both \u2014 geo-tagged markers"),
    ("Results \u2014 Land Cover Classification", "landcover_comparison_peenya.png", "ML-based land cover: 2020 vs 2024 \u2014 clear vegetation-to-built-up transition"),
]:
    img_slide(t, f, c)

# === SLIDES 27-29: Closing ===
bullet_slide("Web Dashboard \u2014 Interactive Monitoring", [
    "Technology: Flask (Python) backend + HTML/JS/CSS frontend",
    "", "Features:",
    "  \u2022 Interactive Leaflet.js map with color-coded violation markers",
    "  \u2022 Zone switching tabs (Peenya / Whitefield)",
    "  \u2022 6 metric cards: Green Cover, Net Change, Veg Loss, Construction, Risk",
    "  \u2022 Violations table with geo-coordinates and area in hectares",
    "  \u2022 ML performance panel with feature importance bars",
    "  \u2022 Tabbed satellite imagery viewer (7 image types)",
    "", "Data Flow:",
    "  \u2022 Pipeline runs OFFLINE (once) \u2192 generates JSON + PNG files",
    "  \u2022 Flask reads pre-computed files \u2192 serves to browser",
    "  \u2022 Zero computation at runtime \u2014 pure file serving",
])

bullet_slide("Scalability & Extensibility", [
    "Geographic Scalability",
    "  \u2022 Pipeline parameterized by bounding box coordinates",
    "  \u2022 Add new zone = add one bbox line + re-run pipeline",
    "  \u2022 KGIS shapefile has all 226 Karnataka taluks",
    "  \u2022 Sentinel-2 freely available globally",
    "", "Temporal Scalability",
    "  \u2022 Sentinel-2 revisits every 5 days",
    "  \u2022 Pipeline can be scheduled (cron) for continuous monitoring",
    "", "Modular Architecture",
    "  \u2022 Each stage is an independent Python module",
    "  \u2022 Any module can be replaced or upgraded independently",
    "  \u2022 preprocess \u2192 indices \u2192 change_detection \u2192 ml_classifier \u2192 visualize",
])

bullet_slide("Conclusion", [
    "1. Built end-to-end automated pipeline for industrial compliance monitoring",
    "",
    "2. Integrated KGIS taluk boundary data for precise study area delineation",
    "",
    "3. Analyzed two zones across 2020\u20132024 using Sentinel-2 satellite imagery",
    "",
    "4. Achieved 97.7\u201398.7% classification accuracy with non-circular RF validation",
    "",
    "5. Detected 126 geo-tagged violation zones (54 Peenya + 72 Whitefield)",
    "",
    "6. Quantified green cover loss: \u221224.1% (Peenya) and \u221227.8% (Whitefield)",
    "",
    "7. Generated automated compliance reports with visual evidence",
])

# === SLIDE 30: Thank You ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
bar(s, Inches(0), Inches(3.0), prs.slide_width, Inches(0.05), ACCENT)
txt(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8), "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.3), Inches(11), Inches(0.6), "Questions & Discussion", 24, c=RGBColor(0x93,0xAE,0xDB), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.5), Inches(11), Inches(0.8), "Future Work: KGIS cadastral integration \u2022 Real-time monitoring \u2022 CNN architectures\nGround truth validation \u2022 All Karnataka zones \u2022 KSPCB system integration", 16, c=RGBColor(0xA0,0xA0,0xC0), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5), "Oishik Kar  \u2022  Mihir Dagar  \u2022  Shouvik Das  \u2022  Youvansh Chauhan", 16, c=RGBColor(0x90,0x90,0xB0), a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.3), Inches(11), Inches(0.5), "IIIT Dharwad  |  ML Course Project  |  Statement 3", 14, c=RGBColor(0x70,0x70,0x90), a=PP_ALIGN.CENTER)

out = os.path.join(BASE, "reports", "presentation.pptx")
prs.save(out)
print(f"PPT saved: {out}")
print(f"Total slides: {len(prs.slides)}")
